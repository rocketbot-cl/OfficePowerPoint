# coding: utf-8
"""
Base para desarrollo de modulos externos.
Para obtener el modulo/Funcion que se esta llamando:
     GetParams("module")

Para obtener las variables enviadas desde formulario/comando Rocketbot:
    var = GetParams(variable)
    Las "variable" se define en forms del archivo package.json

Para modificar la variable de Rocketbot:
    SetVar(Variable_Rocketbot, "dato")

Para obtener una variable de Rocketbot:
    var = GetVar(Variable_Rocketbot)

Para obtener la Opcion seleccionada:
    opcion = GetParams("option")


Para instalar librerias se debe ingresar por terminal a la carpeta "libs"
    
    pip install <package> -t .

"""
import os
import sys

base_path = tmp_global_obj["basepath"]
cur_path = base_path + 'modules' + os.sep + \
    'OfficePowerPoint' + os.sep + 'libs' + os.sep
if cur_path not in sys.path:
    sys.path.append(cur_path)
docto = os.path.join(cur_path.replace("libs", "bin"), "docto.exe")

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx import Presentation


def makeTmpDir(name):
    try:
        os.mkdir("tmp")
        os.mkdir("tmp" + os.sep + name)
    except:
        try:
            os.mkdir("tmp" + os.sep + name)
        except:
            pass

    return os.sep.join(["tmp", name])

def find_shape(shapes, condition, filter_type="shape_id"):
    for shape in shapes:
        print(shape.shape_id, type(shape.shape_id), condition, type(condition))
        if condition == getattr(shape, filter_type):
            print(shape)
            return shape

def change_text_frame(paragraph, text, align="", bold=None, italic=None, underline=None, size=""):
            paragraph.text = text
        
            if size:
                paragraph.font.size = size
            if bold is not None:
                paragraph.font.bold = bold
            if italic is not None:
                paragraph.font.italic = italic
            paragraph.font.underline = underline
            paragraph.alignment = align

ALIGMENT = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }

module = GetParams("module")
global prs
global slide
global slide_layout

try:
    if module == "new":
        prs = Presentation()

    if module == "open":
   
        path = GetParams("path")
        prs = None
        prs = Presentation(path)

    elif module == "get_slides_types":
        res = GetParams("res")
        list_types = []
    
        for type_slide in prs.slide_layouts:
            list_types.append(type_slide.name)
    
        SetVar(res, list_types)


    elif module == "new_slide":
        name = GetParams("type")    

        list_types = []
        for type_slide in prs.slide_layouts:
            print(type_slide.name)
            list_types.append(type_slide.name)
        try:
            name = int(name)
            name = list_types[name]
        except ValueError: 
            pass
        slide_layout = prs.slide_layouts.get_by_name(name)
        slide = prs.slides.add_slide(slide_layout)

    elif module == "save":
        path = GetParams("path")
    
        if path:
            prs.save(path)
    

    elif module == "write":
        text = GetParams("text").replace("\\n", "\n")
        slide_index = GetParams("slide_index")
        shape_index = GetParams("type")
        align_type = GetParams("align")
        size = GetParams("size")
        is_bold = GetParams("bold")
        is_ital = GetParams("italic")
        is_under = GetParams("underline")
        
        align = ALIGMENT["center"]
        if align_type is not None:
            align = ALIGMENT[align_type]
        
        if size:
            size = Pt(int(size))

        bold = ital = under = False
        if is_bold is not None:
            bold = eval(is_bold)
        if is_ital is not None:
            ital = eval(is_ital)
        if is_under is not None:
            under = eval(is_under)

        if not slide_index.isdigit():
            raise Exception("Slide index is not valid")

        slide = prs.slides[int(slide_index)]
        shapes = slide.shapes
        for e in shapes:
            print(e.name)
        if not shape_index.isdigit():
            shape = find_shape(shapes, shape_index, "name")
        else:
            shape = find_shape(shapes, shape_index)
       
        
        p = shape.text_frame.paragraphs[0]
        change_text_frame(p, text, align, bold, ital, under, size)
        

    elif module == "close":
        prs = None

    elif module == "add_pic":

        img_path = GetParams("img_path")
        index = GetParams("slide")
        pos = GetParams("position")
        height = GetParams("height")

        index = int(index)
        if index:
            slide = prs.slides[index]
            slide_layout = slide
        placeholders = slide_layout.shapes.placeholders
        list_types = []
        for type_slide in prs.slide_layouts:
            list_types.append(type_slide.name)
        if slide_layout.name == list_types[8]:
            for placeholder in placeholders:
                if "picture" in placeholder.name.lower():
                    placeholder.insert_picture(img_path)
        else:
            top, left = eval(pos)
            top, left = Pt(float(top)), Pt(float(left))
            if height:
                height = Pt(float(height))
                pic = slide.shapes.add_picture(
                    img_path, left=left, top=top, height=height)
            else:
                pic = slide.shapes.add_picture(img_path, left=left, top=top)

    elif module == "addTextbox":
        text = GetParams("text")
        position = GetParams("position")
        size = GetParams("size")
        index = GetParams("slide")
        align_type = GetParams("align")
       
        is_bold = GetParams("bold")
        is_ital = GetParams("italic")
        is_under = GetParams("underline")
        
        align = ALIGMENT["center"]
        if align_type:
            align = ALIGMENT[align_type]
        

        bold = ital = under = False
        if is_bold is not None:
            bold = eval(is_bold)
        if is_ital is not None:
            ital = eval(is_ital)
        if is_under is not None:
            under = eval(is_under)

        # if not index.isdigit():
        #     raise Exception("Slide index is not valid")

        # slide = prs.slides[int(index)]
        position = position.split(",")
        size = size.split(",")
        pos = position + size
        for i in range(len(pos)):
            pos[i] = Pt(int(pos[i]))

        txt_box = slide.shapes.add_textbox(pos[0], pos[1], pos[2], pos[3])

        p = txt_box.text_frame.paragraphs[0]
        change_text_frame(p, text, align, bold, ital, under)
        txt_box.text_frame.text = text


    elif module == "editText":

        search_by = GetParams("searchBy")
        index = GetParams("slide")
        shape_index = GetParams("shape")
        text = GetParams("text")
        align_type = GetParams("align")
        size = GetParams("size")
        is_bold = GetParams("bold")
        is_ital = GetParams("italic")
        is_under = GetParams("underline")

        align = ALIGMENT["center"]
        if align_type:
            align = ALIGMENT[align_type]
        
        if size:
            size = Pt(int(size))

        bold = ital = under = None
        if is_bold:
            bold = eval(is_bold)
        if is_ital:
            ital = eval(is_ital)
        if is_under:
            under = eval(is_under)

        if not index.isdigit():
            raise Exception("Slide index is not valid")

        slide = prs.slides[int(index)]
        shapes = slide.shapes
        if not shape_index.isdigit():
            shape = find_shape(shapes, shape_index, "name")
        else:
            shape = find_shape(shapes, int(shape_index))

        p = shape.text_frame.paragraphs[0]
        change_text_frame(p, text, align, bold, ital, under, size)


    elif module == "listSlides":
        res = GetParams("res")
        slide_index = int(GetParams("slide_index"))
    
        text_runs = []
        slide = prs.slides[slide_index]
        
        for shape in slide.shapes:
            
            text = ""
            index = shape.shape_id
            name = shape.name
            type_ = shape.shape_type
            if shape.has_text_frame:
                text = shape.text
                
            text_runs.append({
                "id": index,
                "name": name,
                "type": type_,
                "value": text
            })
        SetVar(res, text_runs)

    elif module == "editElement":
        index = GetParams("slide_index")
        shape_index = GetParams("type")
        pos = GetParams("position")
        size = GetParams("size")
        rotation = GetParams("rotation")


        if not index.isdigit():
            raise Exception("Slide index is not valid")

        slide = prs.slides[int(index)]
        shapes = slide.shapes
        if not shape_index.isdigit():
            shape = find_shape(shapes, shape_index, "name")
        else:
            shape = find_shape(shapes, int(shape_index))

        if pos:
            left, top = pos.split(",")
            shape.left = int(left)
            shape.top = int(top)
        if size:
            width, height = size.split(",")
            shape.width = int(width)
            shape.height = int(height)

        if rotation:
            shape.rotation = float(rotation)
        

        print(dir(shape))



except Exception as e:
    print("\x1B[" + "31;40mError\x1B[" + "0m")
    PrintException()
    raise e
